import streamlit as st
import os
import tempfile
from PyPDF2 import PdfReader
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import yaml
import pymongo
import base64
from gridfs import GridFS
from PIL import Image
import hashlib
import requests
import fitz  # PyMuPDF
from pymongo.mongo_client import MongoClient
from pymongo.server_api import ServerApi
import shutil
import pandas as pd
import time
# for similarity check 
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from datetime import datetime
from loginform import fill_login_form


# 設定 Streamlit 網頁為 wide mode(layout="wide") or center(layout="centered")
# st.set_page_config(page_title="Paper Helper", layout="wide")



# streamlit run app_summary.py

# Configure Google Gemini API


# read config from yaml file
file_path = os.path.join(os.path.dirname(__file__), 'config.yaml')
with open(file_path, 'r') as file:
    config = yaml.safe_load(file)

cloud_server_address = config['mongodb_atlas_server_address']
local_server_address = config['local_server_address']


# read api key from yaml file
file_path = os.path.join(os.path.dirname(__file__), 'api_key.yaml')
with open(file_path, 'r') as file:
    api_key = yaml.safe_load(file)

genai.configure(api_key=api_key['Gemini_api_key'])

gemini_model_name = api_key['gemini_model_name']
model = genai.GenerativeModel(model_name=gemini_model_name)



def calculate_similarity(text1, text2):
    vectorizer = TfidfVectorizer().fit_transform([text1, text2])
    vectors = vectorizer.toarray()
    similarity = round(cosine_similarity(vectors)[0, 1],4)
    if similarity > 0.9:
        print(f"文本相似度: {similarity}, 兩篇論文相似")
        same_or_not = True
    else:
        print(f"文本相似度: {similarity}, 兩篇論文不相似")
        same_or_not = False
    return same_or_not,similarity

def paper_similarity_checker(new_paper_text, extracted_data_from_db):   
    # Print the columns of the DataFrame for debugging
    print("Columns in extracted_data_from_db:", extracted_data_from_db.columns.tolist())
    
    for index, row in extracted_data_from_db.iterrows():
        # Check if 'paper_text' exists in the row
        if 'paper_text' not in row:
            print(f"Row {index} does not contain 'paper_text'. Available columns: {row.index.tolist()}")
            continue  # Skip this iteration if 'paper_text' is missing
        
        paper_text = row['paper_text']
        same_or_not, similarity = calculate_similarity(new_paper_text, paper_text)
        print(f"與資料庫中的論文(paper_title: {row['paper_title']})相似度: {similarity}, 是否相似: {same_or_not}")
        if same_or_not:
            return row['paper_title']
    return None
 

def generate_professor_questions(pdf_text, previous_questions):
    """生成教授問題與答案"""
    try:
        question_prompt = (
            f"Based on the following excerpt from the academic paper, generate a thoughtful, concise and well-structured question that reflects the concepts discussed (no more than 20 words):\n\n{pdf_text[:4000]}, question should be different from previous questions, 問法跟主題也要儘量不同: {previous_questions}"
        )
        questions_response = model.generate_content(question_prompt)
        print(questions_response.text.strip())
        answer_prompt = (
            f"Using the Information from the academic paper, provide a comprehensive answer to this question「{questions_response.text.strip()}」(no more than 40 words):\n\n{pdf_text[:4000]}"
        )
        answers_response = model.generate_content(answer_prompt)
        
        question = questions_response.text.strip().split('\n')[0] if questions_response else "No question generated."
        answer = answers_response.text.strip().split('\n')[0] if answers_response else "No answer generated."
        
        return question.strip(), answer.strip()
    except Exception as e:
        st.error(f"Error generating questions and answers: {e}")
        return None, None

def convert_pdf_to_images(pdf_path, output_folder):
    """
    Convert a PDF file to images, saving each page as a separate JPG file.
    
    Args:
        pdf_path (str): Path to the input PDF file
        output_folder (str): Directory where the output images will be saved
    """
    # 如果資料夾不存在，則創建資料夾
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # 獲取 PDF 檔名（不包含副檔名）
    pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]
    print(pdf_name)
    print(pdf_path)
    # 打開 PDF
    doc = fitz.open(pdf_path)

    # 遍歷每一頁
    for i in range(len(doc)):
        page = doc.load_page(i)  # 加載頁面
        pix = page.get_pixmap()  # 將頁面轉為圖片
        output_image_path = os.path.join(output_folder, f'page_{i + 1}.jpg')
        pix.save(output_image_path)  # 儲存圖片
        # print(f"Saved: {output_image_path}")

    doc.close()

def connect_to_mongodb(table_name="PaperHelper", collection_name="file", server_address=cloud_server_address, server_or_atlas="atlas"):
    """
    Connect to MongoDB and return the specified collection.
    
    Parameters:
        table_name (str): The name of the database to connect to.
        collection_name (str): The name of the collection to use.

    Returns:
        pymongo.collection.Collection: The MongoDB collection object.
    """
    try:
        # Connect to the MongoDB server
        if server_or_atlas == "atlas":
            client = MongoClient(server_address, server_api=ServerApi('1')) # for cloud mongodb
        else:
            client = pymongo.MongoClient(server_address) # for local mongodb

        try:
            client.admin.command('ping')
            print("Pinged your deployment. You successfully connected to MongoDB!")
        except Exception as e:
            print(e)
        

        print("Connected to MongoDB successfully.")

        # Select the database (creates it if it doesn't exist)
        db = client[table_name]

        # Check if the database already exists
        if table_name in client.list_database_names():
            print(f"Database '{table_name}' already exists.")
        else:
            print(f"Database '{table_name}' does not exist. It will be created upon adding data.")
        # Select the collection (creates it if it doesn't exist)
        collection = db[collection_name]

        # Check if the collection already exists
        if collection_name in db.list_collection_names():
            print(f"Collection '{collection_name}' already exists.")
        else:
            print(f"Collection '{collection_name}' does not exist. It will be created upon adding data.")

        # Return the collection object for further use
        return collection
    except pymongo.errors.ConnectionError as e:
        print(f"Failed to connect to MongoDB: {e}")
        return None
    except Exception as e:
        print(f"An error occurred: {e}")
        return None

def read_pdf(pdf_path):
    """Read PDF and merge text content from all pages"""
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def generate_paper_title(pdf_text):
    """Generate paper title using Google Gemini API"""
    title_prompt = f"Extract or generate the most appropriate academic paper title from the following text. Only output the title:\n\n{pdf_text[:1000]}"
    try:
        title_response = model.generate_content(title_prompt)
        return title_response.text.strip() if title_response else "Untitled Paper"
    except Exception as e:
        st.error(f"Error generating title: {e}")
        return "Untitled Paper"

def generate_paper_summary(pdf_text):
    """Generate summary and key points using Google Gemini API"""
    summary_prompt = f"Please generate a comprehensive summary of the following PDF text, including main points and key insights. Provide the summary in clear, concise language:\n\n{pdf_text[:2000]}"
    try:
        summary_response = model.generate_content(summary_prompt)
        return summary_response.text.strip() if summary_response else "Unable to generate summary."
    except Exception as e:
        st.error(f"Error generating summary: {e}")
        return "An error occurred while generating the summary."



def generate_pdf_summary_and_details_for_ppt(pdf_text):
    """使用 Google Gemini API 生成 PDF 條列摘要與詳細重點整理"""
    summary_prompt = f"請為以下 PDF 文本生成條列式摘要，每點以「-」開頭，no more than 30 words， 約 ppt 半頁左右內容(用全英文）：\n\n{pdf_text[:2000]}"
    details_prompt = f"請為以下 PDF 文本生成更詳細的重點整理，不少於 500 字(用全英文）：\n\n{pdf_text[:2000]}"
    try:
        summary_response = model.generate_content(summary_prompt)
        summary = summary_response.text.strip() if summary_response else "未能生成摘要。"

        details_response = model.generate_content(details_prompt)
        details = details_response.text.strip() if details_response else "未能生成重點整理。"

        return summary, details
    except Exception as e:
        st.error(f"生成摘要與詳細時發生錯誤：{e}")
        return "未能生成摘要。", "未能生成詳細重點整理。"
 
def generate_image_description_for_page(image_path):
    """使用 Google Gemini API 上傳圖片並生成描述"""
    try:
        myfile = genai.upload_file(image_path)
        prompt = "Can you tell me about the photo? in concise way in Full English, 大概100字左右 ppt 半頁左右內容，用淺想易懂方式解釋"
        result = model.generate_content([myfile, "\n\n", prompt])
        return result.text.strip() if result else f"Unable to generate description for: {image_path}"
    except Exception as e:
        st.error(f"Error generating image description: {e}")
        return "Unable to generate image description"
    

def generate_image_description(image_path):
    """使用 Google Gemini API 上傳圖片並生成描述"""
    try:
        myfile = genai.upload_file(image_path)
        prompt = "Can you tell me about the photo? in concise way in Full English, 大概40字左右 ppt 半頁左右內容"
        result = model.generate_content([myfile, "\n\n", prompt])
        return result.text.strip() if result else f"Unable to generate description for: {image_path}"
    except Exception as e:
        st.error(f"Error generating image description: {e}")
        return "Unable to generate image description"




def add_image_and_description_to_slide(prs, image_path, description, max_chars_per_slide=400):
    """為每張圖片添加解說內容到單獨的幻燈片"""
    description_pages = [description[i:i + max_chars_per_slide] for i in range(0, len(description), max_chars_per_slide)]

    for page_num, page_content in enumerate(description_pages):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # 添加圖片
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(0.5), Inches(8), Inches(4.5))

        # 添加文字框
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(8), Inches(2))
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = text_frame.add_paragraph()
        p.text = page_content.strip()
        p.font.size = Pt(16)
        p.line_spacing = Pt(20)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

def add_text_slide(prs, title, content, bullet=False, font_size=20, max_chars_per_slide=800):
    """添加純文字幻燈片，支持分頁"""
    content_pages = [content[i:i + max_chars_per_slide] for i in range(0, len(content), max_chars_per_slide)]

    for page_num, page_content in enumerate(content_pages):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"{title}（第 {page_num + 1} 部分）"
        text_box = slide.placeholders[1]
        text_frame = text_box.text_frame
        text_frame.clear()

        if bullet:
            for line in page_content.split("\n"):
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(font_size)
        else:
            text_frame.text = page_content
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.alignment = PP_ALIGN.LEFT

def create_presentation(pdf_path, image_paths):
    """根據 PDF 和圖片生成 PPT，返回 PPT 檔案的路徑"""
    pdf_text = read_pdf(pdf_path)
    title = generate_paper_title(pdf_text)
    summary, details = generate_pdf_summary_and_details_for_ppt(pdf_text)
    st.session_state["results"]["paper_title"] = title
    st.session_state["results"]["summary"] = summary
    st.session_state["results"]["pdf_text"] = pdf_text
   
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "由 Paper Helper 生成的內容"

    add_text_slide(prs, "PDF 條列摘要", summary, bullet=True, font_size=20)
    add_text_slide(prs, "PDF 詳細整理", details, bullet=False, font_size=18)

    for image_path in image_paths:
        description = generate_image_description(image_path)
        add_image_and_description_to_slide(prs, image_path, description)

    output_path = os.path.join(tempfile.gettempdir(), f"{title}_presentation.pptx")
    prs.save(output_path)
    return output_path

 

def image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def main():
    if "results" not in st.session_state:
      st.session_state["results"] = {
          "paper_title": None,
          "summary": None,
          "pdf_text": None,
          "image_paths": [],
          "image_descriptions": [],
          "ppt_path": None
       }

    # Custom CSS for styling
    st.markdown("""
    <style>
    .main-title {
        text-align: center;
        font-size: 36px;
        color: #4CAF50;
    }
    .sub-title {
        text-align: center;
        font-size: 18px;
        color: #555;
    }
    .upload-section {
        padding: 20px;
        background: #f9f9f9;
        border-radius: 10px;
        margin-bottom: 20px;
    }
    .result-section {
        background: #f0f0f0;
        border-radius: 10px;
        padding: 20px;
        margin-top: 20px;
    }
    .paper-title {
        text-align: center;
        font-size: 24px;
        color: #333;
        margin-bottom: 15px;
    }
                
    .flip-container {
        display: flex;
        justify-content: center;
        margin-top: 20px;
    }
    .flip-card {
        background-color: transparent;
        width: 500px; /* 擴大寬度 */
        height: 300px; /* 擴大高度 */
        perspective: 1000px;
    }
    .flip-card-inner {
        position: relative;
        width: 100%;
        height: 100%;
        text-align: center; /* 改為居中對齊 */
        transition: transform 0.6s;
        transform-style: preserve-3d;
        box-shadow: 0 4px 8px 0 rgba(0,0,0,0.2);
    }
    .flip-card:hover .flip-card-inner {
        transform: rotateY(180deg);
    }
    .flip-card-front, .flip-card-back {
        position: absolute;
        width: 100%;
        height: 100%;
        backface-visibility: hidden;
        display: flex;
        flex-direction: column; /* 內容縱向排列 */
        justify-content: flex-start; /* 將內容推到卡片的上方 */
        align-items: center; /* 水平居中對齊 */
        border-radius: 10px;
        padding: 20px;
        overflow-wrap: break-word;
        font-size: 18px;
        line-height: 1.5;
    }
    .flip-card-front {
        background-color: #e3f2fd;
        color: black;
    }
    .flip-card-back {
        background-color: #4CAF50;
        color: white;
        transform: rotateY(180deg);
    }
    .question-header, .answer-header {
        font-size: 25px;
        font-weight: bold;
        margin-bottom: 10px;
        text-align: center;
 
    }
    .question-text, .answer-text {
        text-align: center;
        font-size: 20px;
        margin-top: 20px;
        position: absolute;
        top: 50%;
        left: 50%;
        transform: translate(-50%, -50%);
        width: 95%; /* Increased width to reduce margins */
        padding: 0 10px; /* Reduced horizontal padding */
        box-sizing: border-box; /* Ensure padding is included in width */
    }
    .user-history-table {
        font-size: 16px;
        line-height: 1.5;
    }  
    .user-history-title {
        font-size: 24px;
        font-weight: bold;
        margin: 20px 0;
        text-align: center;
        color: #2196F3;
        text-transform: uppercase;
        letter-spacing: 1px;
        text-shadow: 1px 1px 2px rgba(0,0,0,0.1);
    }
    .stDeployButton {
        visibility: hidden;
    }
                
    .logout-button {
        position: fixed;
        top: 10px;
        right: 10px;
        z-index: 1000;
    }
    </style>
    """, unsafe_allow_html=True)

    # reading the comfig file
    print("start to read the config file")
    config_file_path = os.path.join(os.path.dirname(__file__), 'config.yaml')
    with open(config_file_path, 'r') as file:
        config = yaml.safe_load(file)

    pdf_to_image_folder = config['pdf_to_image_folder']
    cropped_image_folder = config['cropped_image_folder']
    logout_page_url = config['logout_page_url']
    print("done read the config file\n")
 

    ## 初始化：清空pdf_to_image_folder資料夾
    if os.path.exists(pdf_to_image_folder):
        shutil.rmtree(pdf_to_image_folder)
        # create the folder
    else:
        os.makedirs(pdf_to_image_folder, exist_ok=True)

    # App title and description
    st.markdown('<h1 class="main-title">Paper Helper</h1>', unsafe_allow_html=True)
    st.markdown('<p class="sub-title">Upload your PDF and images to get comprehensive insights</p>', unsafe_allow_html=True)
    
    # User past generated history - section
    st.markdown('<div class="upload-section">', unsafe_allow_html=True)

    if "df_extracted_data" not in st.session_state:
        with st.spinner("Loading DB..."):
            # Connect to database and get data
            file_collection = connect_to_mongodb(
                table_name="PaperHelper", 
                collection_name="file", 
                server_address=cloud_server_address, 
                server_or_atlas="atlas")
            data = file_collection.find()
            

        # 要取用的欄位    
        selected_columns = ["user_id","generate_time","paper_title", "paper_summary", "paper_text"]


        all_docs = []
        for doc in data:
            # Access each column individually and create a dict with selected fields
            selected_data = {}
            for col in selected_columns:
                selected_data[col] = doc.get(col, "") # Use get() with default empty string if key not found
            all_docs.append(selected_data)
            
        # Create single dataframe from all documents
        df_extracted_data = pd.DataFrame(all_docs)
        df_extracted_data = df_extracted_data.sort_values(by="generate_time", ascending=False)
        # save in session state
        st.session_state["df_extracted_data"] = df_extracted_data
    else:
        df_extracted_data = pd.DataFrame(st.session_state["df_extracted_data"]) # make sure it is a dataframe
        print("df_extracted_data", df_extracted_data)
        print("type of df_extracted_data", type(df_extracted_data))
        print("columns of df_extracted_data", df_extracted_data.columns)
        
 
        

    st.markdown('<p class="user-history-title">User History</p>', unsafe_allow_html=True)
    df_user_history = df_extracted_data.rename(columns={"user_id": "User ID", "generate_time": "Generate Time", "paper_title": "Paper Title", "paper_summary": "Paper Summary"})
    user_history_columns = ["User ID", "Generate Time", "Paper Title", "Paper Summary"]
    st.dataframe(df_user_history[user_history_columns])
            
    st.markdown('</div>', unsafe_allow_html=True)




    # Upload section
    st.markdown('<div class="upload-section">', unsafe_allow_html=True)
    uploaded_pdf = st.file_uploader("Choose a PDF file", type=["pdf"])
    print("done upload the pdf file\n")

    # 上傳圖片 => 改成從特定資料夾讀取
    # uploaded_images = st.file_uploader("Choose images (multiple allowed)", type=["png", "jpg", "jpeg"], accept_multiple_files=True)

    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Session state variables to retain results
    if "results" not in st.session_state:
        st.session_state["results"] = {}


    # if uploaded_pdf => start to convert pdf to images
    if uploaded_pdf:
        # 將 Streamlit 的 UploadedFile 儲存為臨時檔案
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_pdf.write(uploaded_pdf.read())
            temp_pdf_path = temp_pdf.name
        uploaded_pdf.seek(0)  # 重置指標，讓檔案可以再次被讀取

        # 轉換pdf成圖片
        convert_pdf_to_images(
            temp_pdf_path,
            output_folder=pdf_to_image_folder
        )
        print("done convert pdf to images\n")


         ## debug for cropped_image_folder
        uploaded_images_debug = os.listdir(cropped_image_folder)
        print('Before clear the cropped_image_folder, the folder contains the images:', uploaded_images_debug)
        for uploaded_image in uploaded_images_debug:
            # Construct the full path to the image
            image_full_path = os.path.join(cropped_image_folder, uploaded_image)
            print(f"Before clear the cropped_image_folder, the folder contains the image: {image_full_path}")

        print("\n--------------------------------\n")




    # add a button to open the cropper image page
    open_cropper_image_button = st.button("Image Cropper ")
    # 如果按下cropper image button => 開啟裁切器 (try-project/express/public/index.html)
    if open_cropper_image_button: 
        st.markdown('<iframe src="https://zy.papperhelper.xyz" width="120%" height="1100px"></iframe>', unsafe_allow_html=True)





    generate_insights_and_slides_button = st.button("Generate Insights, Slides and Questions")
    if st.button("Logout"):
        st.session_state.logged_in = False
        st.success("You have successfully logged out!")
        
        # Redirect to the login page using JavaScript
        login_page = logout_page_url
        st.markdown(f'<meta http-equiv="refresh" content="0; url={login_page}">', unsafe_allow_html=True)


    if generate_insights_and_slides_button != True:
        # only clear when the pdf is uploaded, and
        # clear image in the cropped_image_folder, but not the folder itself
        print("start to clear the cropped_image_folder")
        image_list = os.listdir(cropped_image_folder)
        for file in image_list:
            os.remove(os.path.join(cropped_image_folder, file))
        print("done clear the cropped_image_folder\n")


    uploaded_images = os.listdir(cropped_image_folder)

    ## debug
    print('uploaded_images', uploaded_images)
    for uploaded_image in uploaded_images:
        # Construct the full path to the image
        image_full_path = os.path.join(cropped_image_folder, uploaded_image)
        print(f"folder contains the image: {image_full_path}")

    print("\n--------------------------------\n")

    pdf_upload = config['pdf_upload_path']


    if uploaded_pdf:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf", dir=pdf_upload) as temp_pdf:
            temp_pdf.write(uploaded_pdf.read())
            temp_pdf_path = temp_pdf.name
        uploaded_pdf.seek(0)  # 重置指標，讓檔案可以再次被讀取

        # Read PDF content and generate title and summary
        pdf_text = read_pdf(temp_pdf_path)
        st.session_state["results"]["pdf_text"] = pdf_text

        # 透過 alculate_similarity 計算new paper 與 df 中的paper_text 論文相似度
        paper_existed_in_db_or_not = paper_similarity_checker(new_paper_text=st.session_state["results"]["pdf_text"], extracted_data_from_db=df_extracted_data)
        print("paper_existed_in_db_or_not", paper_existed_in_db_or_not)
        if paper_existed_in_db_or_not != None:
            paper_title = paper_existed_in_db_or_not

            # 搜索符合條件的行，並提取 paper_summary 的內容
            paper_summary_series = df_extracted_data [df_extracted_data['paper_title'] == paper_title]['paper_summary'].head(1)

            # 如果找到匹配的行，提取文字，否則返回空字串
            paper_summary = paper_summary_series.iloc[0] if not paper_summary_series.empty else "No summary found."

            # 打印 paper_summary 的文字
            print(paper_title, paper_summary)
            # show it on the screen
            st.success("Paper already exists in the database")
            if generate_insights_and_slides_button == False:
                st.markdown(f'<div class="paper-title">Paper Title: {paper_title}</div>', unsafe_allow_html=True)
                st.markdown(f'<div class="paper-summary">Summary: {paper_summary}</div>', unsafe_allow_html=True)
            # Store results in session state
            st.session_state["results"]["paper_title"] = paper_title
            st.session_state["results"]["summary"] = paper_summary




    if uploaded_pdf and generate_insights_and_slides_button:
        print("start to generate summary")
        with st.spinner("Analyzing your document and images..."):
            progress_bar = st.progress(0, "Processing...")
            # Save PDF to a temporary file
            # with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            #     temp_pdf.write(uploaded_pdf.read())
            #     temp_pdf_path = temp_pdf.name

            # Read PDF content and generate title and summary
            pdf_text = read_pdf(temp_pdf_path)
            st.session_state["results"]["pdf_text"] = pdf_text

            if st.session_state["results"]["paper_title"] == None:
                progress_bar.progress(0.1)
                paper_title = generate_paper_title(st.session_state["results"]["pdf_text"])
                progress_bar.progress(0.2)
                summary = generate_paper_summary(st.session_state["results"]["pdf_text"])
                progress_bar.progress(0.4)
            else:
                progress_bar.progress(0.1)
                paper_title = st.session_state["results"]["paper_title"]
                progress_bar.progress(0.2)
                time.sleep(0.1)
                summary = st.session_state["results"]["summary"]
                progress_bar.progress(0.4)
                time.sleep(0.1)

            # Store results in session state
            # st.session_state["results"]["paper_title"] = paper_title
            # st.session_state["results"]["summary"] = summary
            

            print("\n--------------------------------\n")
            print("done generate summary, start to generate image descriptions")

            # Generate image descriptions
            image_paths = []
            image_descriptions = []
            # used images from the cropped_image_folder, instead of st.file_uploader
            uploaded_images = os.listdir(cropped_image_folder)
            ## debug
            print('uploaded_images', uploaded_images)
            for uploaded_image in uploaded_images:
                # Construct the full path to the image
                image_full_path = os.path.join(cropped_image_folder, uploaded_image)
                print(f"folder contains the image: {image_full_path}")

            # Loop through the images in the folder
            for uploaded_image in uploaded_images:
                # Construct the full path to the image
                image_full_path = os.path.join(cropped_image_folder, uploaded_image)
                print(f"Processing: {image_full_path}")

                # 確認檔案存在且不是空檔案
                if not os.path.isfile(image_full_path):
                    print(f"Skipped: {image_full_path} does not exist.")
                    continue
                if os.path.getsize(image_full_path) == 0:
                    print(f"Skipped empty file: {image_full_path}")
                    continue

                try:
                    # 開啟圖片檔案
                    with Image.open(image_full_path) as img:
                        # Convert the image to RGB to ensure compatibility with JPEG format
                        img = img.convert("RGB")
                        print(f"image is opened: {image_full_path}")
                        # 儲存到暫存檔案中
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_image:
                            img.save(temp_image.name, format="JPEG")
                            image_path = temp_image.name  # 暫存檔案的路徑
                            image_paths.append(image_path)

                            # 生成描述並新增到清單
                            description = generate_image_description_for_page(image_path)
                            image_descriptions.append(description)

                except Exception as e:
                    # 如果出現問題，記錄錯誤並跳過
                    print(f"Skipped {uploaded_image} due to an error: {e}")

 

            st.session_state["results"]["image_paths"] = image_paths
            st.session_state["results"]["image_descriptions"] = image_descriptions
            # print(st.session_state["results"]["image_paths"])
            # print(st.session_state["results"]["image_descriptions"])

            # Generate presentation
            ppt_path = create_presentation(temp_pdf_path, image_paths)
            st.session_state["results"]["ppt_path"] = ppt_path
            progress_bar.progress(0.6)
            # connect to mongodb
            file_collection = connect_to_mongodb(table_name="PaperHelper", collection_name="file", server_address=cloud_server_address, server_or_atlas="atlas")
            # Initialize GridFS
            fs = GridFS(file_collection.database)

            # Generate a common ID based on the paper title
            paper_title = st.session_state["results"]["paper_title"]
            common_id = hashlib.md5(paper_title.encode()).hexdigest()

            # Combine all data to be saved in a single dictionary
            data_to_save = {
                "common_id": common_id,
                "user_id": None,
                "generate_time": datetime.now().strftime("%Y-%m-%d %H:%M"),
                "paper_text": st.session_state["results"]["pdf_text"],
                "paper_title": paper_title,
                "paper_summary": st.session_state["results"]["summary"],
                "image_paths": st.session_state["results"].get("image_paths", []),
                "image_descriptions": st.session_state["results"].get("image_descriptions", []),
                "images_base64": []
            }

            # Transform all images to base64 and add to the dictionary
            for image_path in data_to_save["image_paths"]:
                image_base64 = image_to_base64(image_path)
                data_to_save["images_base64"].append(image_base64)

            # Read the ppt file and save it to GridFS
            if "ppt_path" in st.session_state["results"]:
                with open(st.session_state["results"]["ppt_path"], "rb") as ppt_file:
                    ppt_file_id = fs.put(ppt_file, filename=os.path.basename(st.session_state["results"]["ppt_path"]))
                    data_to_save["ppt_file_id"] = ppt_file_id

            # Save the combined data to MongoDB
            file_collection.insert_one(data_to_save)

            # Display results if available in session state
            if "results" in st.session_state and "paper_title" in st.session_state["results"]:
                st.markdown('<div class="result-section">', unsafe_allow_html=True)
                st.markdown(f'<div class="paper-title" style="font-size: 28px;">Paper Title: {st.session_state["results"]["paper_title"]}</div>', unsafe_allow_html=True)
                st.subheader("Paper Summary")
                st.write(st.session_state["results"]["summary"])
                st.markdown('</div>', unsafe_allow_html=True)

                if "image_descriptions" in st.session_state["results"]:
                    st.markdown('<div class="result-section">', unsafe_allow_html=True)
                    st.subheader("Image Descriptions")
                    for i, (image_path, description) in enumerate(zip(st.session_state["results"]["image_paths"], st.session_state["results"]["image_descriptions"])):
                        st.markdown(f"### Image {i + 1}")
                        st.image(image_path, width=800)
                        st.write(description)
                    st.markdown('</div>', unsafe_allow_html=True)
 
            
    
        with st.spinner("Analyzing your document and generating questions..."):
            # Save PDF to a temporary file
            pdf_text = st.session_state["results"]["pdf_text"]
            num_of_questionss = 3
            previous_questions = []
            st.markdown(f"### 👨‍🏫 Q&A Section")
            for i in range(num_of_questionss): 
                question, answer = generate_professor_questions(pdf_text, previous_questions)
                previous_questions.append(question)
                # Display professor question and answer
                if question and answer:
                    
                    st.markdown('<div class="flip-container">', unsafe_allow_html=True)

                    # 僅顯示第一個問題
                    st.markdown(f"""
                    <div class="flip-card">
                        <div class="flip-card-inner">
                            <div class="flip-card-front">
                                <div class="question-header">Question - {i+1}</div>
                                <p class="question-text">{question}</p>
                            </div>
                            <div class="flip-card-back">
                                <div class="answer-header">Answer</div>
                                <p class="answer-text">{answer}</p>
                            </div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)

                    st.markdown('</div>', unsafe_allow_html=True)
                    # 空行
                    st.markdown('<br>', unsafe_allow_html=True)
                    progress_bar.progress(0.8)

        progress_bar.progress(1)
        # delete the progress bar
        progress_bar.empty()
        
        if "ppt_path" in st.session_state["results"]:
            st.success("Slides generated successfully 🎉")
            with open(st.session_state["results"]["ppt_path"], "rb") as ppt_file:
                st.download_button(
                    label="Download Slides",
                    data=ppt_file,
                    file_name=os.path.basename(st.session_state["results"]["ppt_path"]),
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"  
                     
                )


if __name__ == "__main__":
    main()
